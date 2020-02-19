#!/usr/bin/python3

# Copyright (c) 2014-2020, University of Cantabria
# All rights reserved.
#
# Redistribution and use in source and binary forms, with or without
# modification, are permitted provided that the following conditions are met:
#
# Redistributions of source code must retain the above copyright notice, this
# list of conditions and the following disclaimer.
# Redistributions in binary form must reproduce the above copyright notice,
# this list of conditions and the following disclaimer in the documentation
# and/or other materials provided with the distribution.
#
# THIS SOFTWARE IS PROVIDED BY THE COPYRIGHT HOLDERS AND CONTRIBUTORS "AS IS"
# AND ANY EXPRESS OR IMPLIED WARRANTIES, INCLUDING, BUT NOT LIMITED TO, THE
# IMPLIED WARRANTIES OF MERCHANTABILITY AND FITNESS FOR A PARTICULAR PURPOSE
# ARE DISCLAIMED. IN NO EVENT SHALL THE COPYRIGHT OWNER OR CONTRIBUTORS BE
# LIABLE FOR ANY DIRECT, INDIRECT, INCIDENTAL, SPECIAL, EXEMPLARY, OR
# CONSEQUENTIAL DAMAGES (INCLUDING, BUT NOT LIMITED TO, PROCUREMENT OF
# SUBSTITUTE GOODS OR SERVICES; LOSS OF USE, DATA, OR PROFITS; OR BUSINESS
# INTERRUPTION) HOWEVER CAUSED AND ON ANY THEORY OF LIABILITY, WHETHER IN
# CONTRACT, STRICT LIABILITY, OR TORT (INCLUDING NEGLIGENCE OR OTHERWISE)
# ARISING IN ANY WAY OUT OF THE USE OF THIS SOFTWARE, EVEN IF ADVISED OF THE
# POSSIBILITY OF SUCH DAMAGE.

# Author: Ivan Perez

# NOTE: Obsolete script, see ../plotter/ instead. I haven't removed
# this plotter because it generates a PowerPoint file with the Excel's graphs.
# Generates Y-X scattered plots (e.g. latency vs offered load) form
# experiments JSON files
# NOTE: Take a look to the "experiment" directory to see some examples
# NOTE: to use this script you must add your user info in "host.py"

# @TODO: Usage: ./plotter.py <experiment.json>

import matplotlib.pyplot as plt
import pandas as pd
import sys
import os
import json
from pptx import Presentation
from pptx.chart.data import ChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE 
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_LEGEND_POSITION

#viridis_pallete = [[68,1,84],[71,45,123],[59,82,139],[44,114,142],[33,144,140],[39,173,129],[93,200,99],[170,220,50],[254,231,57]] #RGB
viridis_pallete = ["440154","472D7B","3B528B","2C728E","21908C","27AD81","5DC863","AADC32","FEE739"] #RGB

class Chart():
    def __init__(self,experiment):
        with open(experiment,'r') as f:
           data = json.load(f)
        self.data = data
        
class LineChart(Chart):
    def __init__(self, experiment):
        Chart.__init__(self,experiment)
        self.title = self.data["plots"]["title"]
        self.legend = self.data["plots"]["legend"]
        self.xaxis_title = self.data["plots"]["xaxis-title"]
        self.yaxis_title = self.data["plots"]["yaxis-title"]

        if self.data["plots"]["custom"] == "true": 
            self.legend = []
            self.load = []
            self.yaxis = []
            for indx, sim in enumerate(self.data["plots"]["simulations"]):
                br = BookSimReader(sim)
                cl = int(self.data["plots"]["classes"][indx])
                for c in range(cl):
                    self.legend.append(self.data["plots"]["legend"][indx] + " class: " + str(c))
                    self.load.append(br.read_load(cl))
                    self.yaxis.append(br.read_stat(self.data["plots"]["yaxis"],c))
        else:
            print("Generalized version not implemented yet")
    
    def generate_matplotlib(self):
        print("Not implemented yet")

    def generate_power_point(self):
        prs = Presentation(self.data["plots"]["template-pptx"])
        slide = prs.slides.add_slide(prs.slide_layouts[5]) 
        shapes = slide.shapes
        shapes.title.text = self.title

        chart_data = XyChartData()
        
        series = []
        for index, l in enumerate(self.yaxis):
            series.append(chart_data.add_series(self.legend[index]))
            for element in range(len(self.load[index])):
                series[-1].add_data_point(self.load[index][element], self.yaxis[index][element])

        x, y, cx, cy = Inches(2), Inches(2), Inches(9), Inches(4.5) 

        graphic_frame = slide.shapes.add_chart(
                    XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data
                    )
        chart = graphic_frame.chart
#        chart.legend.include_in_layout = True
        chart.legend.position = XL_LEGEND_POSITION.TOP
        chart.category_axis.has_title = True
        chart.category_axis.axis_title.text_frame.text = self.xaxis_title
        chart.category_axis.major_gridline = True
        chart.value_axis.has_title = True
        chart.value_axis.axis_title.text_frame.text = self.yaxis_title
        if "y-max" in self.data["plots"]:
            chart.value_axis.maximum_scale = float(self.data["plots"]["y-max"])

        stride = len(viridis_pallete)/len(chart.series)
        for index,serie in enumerate(chart.series):
            serie.marker.format.fill.background()
            serie.marker.format.line.fill.background()
            serie.format.line.fill.solid()
            serie.format.line.color.rgb = RGBColor.from_string(viridis_pallete[index*stride])

        prs.save(self.data["plots"]["final-pptx"])
        


def main():
    lc = LineChart(sys.argv[1]) 
    lc.generate_power_point()

    #with open(sys.argv[1],'r') as f:
    #    data = json.load(f)

    #    if data["plots"]["type"] == "single_yaxis":
    #        single_yaxis(data)
    #    elif data["plots"]["type"] == "single_yaxis_single_line":
    #        single_yaxis_single_line(data)
    #    elif data["plots"]["type"] == "double_yaxis":
    #        double_yaxis(data)
    #    elif data["plots"]["type"] == "double_yaxis_single_line":
    #        double_yaxis_single_line(data)
    #    elif data["plots"]["type"] == "single_yaxis_single_line_multi_class":
    #        single_yaxis_single_line_multi_class(data)
    #    else:
    #        sys.exit("Error plot type: " + data["plots"]["type"] + " uknown")

def single_yaxis_single_line(data):
    subplots = data["plots"]["subplots"]
    xaxis_stat = data["plots"]["xaxis"]
    yaxis_stat = data["plots"]["yaxis1"]

    for sp in data["variables"][subplots]:
        fig, ax = plt.subplots()
        for f in os.listdir(data["output_dir"]):
            if data["name"] in f and sp in f and ".csv" in f:
                br = BookSimReader(data["output_dir"] + "/" + f)
                ax.plot(br.read_stat(xaxis_stat,-1), br.read_stat(yaxis_stat,-1))


        fig.suptitle(data['name'] + "_" + sp)
        ax.legend(loc=0)
        ax.set_xlabel(data["plots"]["xlabel"])
        ax.set_ylabel(data["plots"]["ylabel1"])
        ax.set_ylim(ymin= float(data["plots"]["ymin1"]), ymax=float(data["plots"]["ymax1"]))
        ax.grid()
        fig.savefig(data["name"] + "_" + sp + ".pdf")

def single_yaxis(data):
    subplots = data["plots"]["subplots"]
    lines = data["plots"]["lines"]
    xaxis_stat = data["plots"]["xaxis"]
    yaxis_stat = data["plots"]["yaxis1"]

    for sp in data["variables"][subplots]:
        fig, ax = plt.subplots()
        for line in data["variables"][lines]:
            for f in os.listdir(data["output_dir"]):
                if data["name"] in f and sp in f and line in f and ".csv" in f:
                    br = BookSimReader(data["output_dir"] + "/" + f)
                    ax.plot(br.read_stat(xaxis_stat,-1), br.read_stat(yaxis_stat,-1), label = lines + " = " + line)


        fig.suptitle(data['name'] + "_" + sp)
        ax.legend(loc=0)
        ax.set_xlabel(data["plots"]["xlabel"])
        ax.set_ylabel(data["plots"]["ylabel1"])
        ax.set_ylim(ymin= float(data["plots"]["ymin1"]), ymax=float(data["plots"]["ymax1"]))
        ax.grid()
        fig.savefig(data["name"] + "_" + sp + ".pdf")

def double_yaxis_single_line(data):
    subplots = data["plots"]["subplots"]
    xaxis_stat = data["plots"]["xaxis"]
    y1axis_stat = data["plots"]["yaxis1"]
    y2axis_stat = data["plots"]["yaxis2"]

    for sp in data["variables"][subplots]:
        fig, ax = plt.subplots()
        for f in os.listdir(data["output_dir"]):
            if data["name"] in f and sp in f and ".csv" in f:
                br = BookSimReader(data["output_dir"] + "/" + f)
                ax.plot(br.read_stat(xaxis_stat,-1), br.read_stat(y1axis_stat,-1), label = y1axis_stat)
                ax.plot(br.read_stat(xaxis_stat,-1), br.read_stat(y2axis_stat,-1), label = y2axis_stat)


        fig.suptitle(data['name'] + "_" + sp)
        ax.legend(loc=0)
        ax.set_xlabel(data["plots"]["xlabel"])
        ax.set_ylabel(data["plots"]["ylabel1"])
        ax.set_ylim(ymin= float(data["plots"]["ymin1"]), ymax=float(data["plots"]["ymax1"]))
        ax.grid()
        fig.savefig(data["name"] + "_" + sp + ".pdf")


def double_yaxis(data):
    subplots = data["plots"]["subplots"]
    lines = data["plots"]["lines"]
    xaxis_stat = data["plots"]["xaxis"]
    y1axis_stat = data["plots"]["yaxis1"]
    y2axis_stat = data["plots"]["yaxis2"]

    asd = 0
    for sp in data["variables"][subplots]:
        fig, ax = plt.subplots()
        for line in data["variables"][lines]:
            for f in os.listdir(data["output_dir"]):
                if data["name"] in f and (lines + "_" + line) in f and ".csv" in f:
                    asd = asd+1
                    br = BookSimReader(data["output_dir"] + "/" + f)
                    ax.plot(br.read_stat(xaxis_stat,-1), br.read_stat(y1axis_stat,-1), label = y1axis_stat + ": " + lines + " = " + line)
                    ax.plot(br.read_stat(xaxis_stat,-1), br.read_stat(y2axis_stat,-1), label = y2axis_stat + ": " + lines + " = " + line)


        fig.suptitle(data['name'] + "_" + sp)
        ax.legend(loc=0)
        ax.set_xlabel(data["plots"]["xlabel"])
        ax.set_ylabel(data["plots"]["ylabel1"])
        ax.set_ylim(ymin= float(data["plots"]["ymin1"]), ymax=float(data["plots"]["ymax1"]))
        ax.grid()
        fig.savefig(data["name"] + "_" + sp + ".pdf")

def single_yaxis_single_line_multi_class(data):
    subplots = data["plots"]["subplots"]
    xaxis_stat = data["plots"]["xaxis"]
    yaxis_stat = data["plots"]["yaxis1"]

    classes = int(data["plots"]["classes"])

    for sp in data["variables"][subplots]:
        fig, ax = plt.subplots()
        for f in os.listdir(data["output_dir"]):
            if data["name"] in f and sp in f and ".csv" in f:
                br = booksimreader(data["output_dir"] + "/" + f)
                for c in range(classes):
                    x_list = br.read_stat(xaxis_stat,c)
                    y_list =  br.read_stat(yaxis_stat,c)
                    if "colors" in data["plots"]:
                        ax.plot(x_list, y_list, label = data["plots"]["legend"][c], color = data["plots"]["colors"][c])
                    else:
                        ax.plot(x_list, y_list, label = data["plots"]["legend"][c])



        fig.suptitle(data['plots']["title"] + sp)
        ax.legend(loc=0)
        ax.set_xlabel(data["plots"]["xlabel"])
        ax.set_ylabel(data["plots"]["ylabel1"])
        ax.set_ylim(ymin= float(data["plots"]["ymin1"]), ymax=float(data["plots"]["ymax1"]))
#        ax.set_xlim(xmin= float(data["plots"]["xmin"]), xmax=float(data["plots"]["xmax"]))
        ax.grid()
#        fig.savefig(data["name"] + "_" + sp + ".pdf")
        fig.savefig(data["name"] + "_" + sp + ".png")


class BookSimReader(object):
    def __init__(self, filename):
        self.filename = filename
        self.datatable = self.create_datatable(filename)
        #self.initial_time = self.datatable.time[0]

    def create_datatable(self, filename):
        dt = pd.read_csv(filename, sep=',')
        return dt

    def read_time(self):
        return self.datatable.loc[self.datatable['class'] == 0]["time"] - self.initial_time

    def read_load(self,cl=0,total=True):
        if total:
            for c in range(cl):

                if c==0:
                    load = self.datatable.loc[self.datatable['class'] == c]['load'].tolist()
                    packet_size = self.datatable.loc[self.datatable['class'] == c]['avg_sent_packet_size'].tolist()
                    a = [x*y for x,y in zip(load,packet_size)]
                else:
                    load = self.datatable.loc[self.datatable['class'] == c]['load'].tolist()
                    packet_size = self.datatable.loc[self.datatable['class'] == c]['avg_sent_packet_size'].tolist()
                    b = [x*y for x,y in zip(load,packet_size)]
                    a = [x+y for x,y in zip(a,b)]
            return a
        else:
            load = self.datatable.loc[self.datatable['class'] == cl]['load']
            packet_size = self.datatable.loc[self.datatable['class'] == c]['avg_sent_packet_size'].tolist()
            return [x*y for x,y in zip(load,packet_size)]

    def read_stat(self, stat, class_id):
        # if class_id is -1 it returns the sum of the stat for all the classes
        if class_id == -1:
            return self.datatable.groupby('load')[stat].sum().tolist()
        else:
            return self.datatable.loc[self.datatable['class'] == class_id][stat].tolist()

    def read_plat_ponderated(self):
        group = self.datatable.groupby("time")["avg_plat","sent_packets"]
        return group.apply(plat_ponderated)

    def read_nlat_ponderated(self):
        group = self.datatable.groupby("time")["avg_nlat","sent_packets"]
        return group.apply(nlat_ponderated)


def plat_ponderated(group):
    stat = group["avg_plat"]
    multiplier = group["sent_packets"]

    return (stat*multiplier).sum() / multiplier.sum()

def nlat_ponderated(group):
    stat = group["avg_nlat"]
    multiplier = group["sent_packets"]

    return (stat*multiplier).sum() / multiplier.sum()



if __name__ == '__main__':
    main()
